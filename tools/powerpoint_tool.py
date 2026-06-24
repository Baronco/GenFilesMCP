import os
from io import BytesIO
from typing import Annotated, List, Literal, Optional, Union

import pandas as pd
import yaml
from pydantic import BaseModel, Field, model_validator
from pptx import Presentation
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt
from utils.charts import chart
from PIL import Image as PILImage
from utils.download_file import download_file
from utils.upload_file import upload_file
from utils.knowledge import create_knowledge
from utils.get_user_id import get_user_id
from utils.authorization import _get_bearer_token
from json import dumps
from utils.logger import get_logger

logger = get_logger(__name__)

# Accepts theme tokens or an explicit hex color (#RRGGBB)
BackgroundField = Union[
    Literal["accent_color", "background_color"],
    Annotated[str, Field(pattern=r"^#[0-9A-Fa-f]{6}$")],
]

TextStyle = Literal["prose", "bullets"]

def generate_powerpoint(python_script, file_name, images_list, request, URL, ENABLE_CREATE_KNOWLEDGE, knowledge_name): 
    """
    Generate a PowerPoint file using an AI-generated Python script.

    Returns:
        dict: Contains 'file_path_download' with a markdown hyperlink for downloading the generated PowerPoint file.
    """
    try:
        images = [] # to save bytesIO images

        if len(images_list) > 0:
            logger.info(f"Received {len(images_list)} images for PPTX generation.")

        # download images
        for idx, image in enumerate(images_list):
            image_file = download_file(URL, _get_bearer_token(request), image)
            if isinstance(image_file, dict) and "error" in image_file:
                return {"error": {"message": f"Error downloading image with ID {image}: {image_file['error']['message']}"}}
            images.append(image_file)
        
        # Prepare the execution context with images
        buffer = BytesIO()
        buffer.name = f'{file_name}.pptx'
        context = {"pptx_buffer": buffer, "images": images}
        try:
            exec(python_script, context)
        except Exception as exec_e:
            return {"error": {"message": f"Error executing script: {str(exec_e)}"}}

        buffer.seek(0)

        try:
            bearer_token = _get_bearer_token(request)
            logger.info("=> Received authorization header!")
        except:
            logger.error("=> Error retrieving authorization header")

        # resolve user_id from token
        user_id = get_user_id(URL, bearer_token)
        if not user_id:
            logger.error("=> Error obtaining user id from token")
            return dumps({"error": {"message": "Error obtaining user id from token"}}, indent=4, ensure_ascii=False)

        response, request_data = upload_file(
            url=URL, 
            token=bearer_token, 
            file_data=buffer,
            filename=file_name,
            file_type="pptx"
        )

        if "file_path_download" in response and ENABLE_CREATE_KNOWLEDGE:
            create_knowledge_status = create_knowledge(
                url=URL,
                token=bearer_token,
                file_id=request_data['id'],
                user_id=user_id,
                knowledge_name=knowledge_name
            )
            if create_knowledge_status:
                logger.info("=> Knowledge base updated successfully.")
            else:
                logger.error("=> Error creating or updating knowledge base")
        elif "error" in response:
            logger.error("=> Error uploading the file.")
        else:
            logger.info("=> Skipping knowledge creation because ENABLE_CREATE_KNOWLEDGE is false")

        return response 
    
    except Exception as e:
        logger.error("=> An unknown error occurred during .PPTX document generation.")
        return dumps(
            {
                "error": {
                    "message": str(e)
                }
            }, 
            indent=4, 
            ensure_ascii=False
        )


# ─────────────────────────── PPTX YAML helpers ────────────────────────────

class Theme(BaseModel):
    """A named, curated bundle of visual choices applied to an entire presentation.
    Replaces the four required `global.*` fields the LLM used to set individually.
    `gradient_accent` is the second stop of the gradient used on "impact" slides
    (cover, section_divider, stat_highlight) — ordinary content slides never use it."""
    accent_color: str
    gradient_accent: str
    background_color: str
    font_heading: str
    font_body: str
    chart_palette: str


def _derive_palette(base_hex: str, mode: str):
    """Derive a coherent (accent, gradient_accent, background) triple from a SINGLE base
    hue so every surface in a theme stays in the same color family — accents, the impact
    gradient, and the slide background are all the same hue at different depths, never a
    mix of unrelated blues/grays. `mode` selects the contrast structure:
      - 'light': near-white background, medium accent, mid-bright gradient end
      - 'dark' : deep clearly-hued background, bright accent, mid-deep gradient end
        (no theme currently uses 'dark' — backgrounds are all light per user preference —
         but the structure is kept so a dark theme can be added later)
    """
    import colorsys
    v = base_hex.lstrip("#")
    r, g, b = int(v[0:2], 16) / 255, int(v[2:4], 16) / 255, int(v[4:6], 16) / 255
    h, _, s = colorsys.rgb_to_hls(r, g, b)
    # Tame chroma so palettes read as refined, not glaring: enough color to be clear, capped
    # so nothing is neon-intense (the user asked for slightly softer, lighter tones).
    s = min(max(s, 0.42), 0.62)

    def hx(light: float, sat: float) -> str:
        rr, gg, bb = colorsys.hls_to_rgb(h, max(0.0, min(1.0, light)), max(0.0, min(1.0, sat)))
        return "#{:02X}{:02X}{:02X}".format(int(rr * 255), int(gg * 255), int(bb * 255))

    if mode == "dark":
        return hx(0.64, s), hx(0.42, s), hx(0.12, min(s * 0.12, 0.07))
    # light mode: medium (not deep) accent + lighter gradient end so covers/headers feel
    # softer; near-white tinted background. Header bars auto-deepen for legible white text.
    return hx(0.38, s), hx(0.60, s), hx(0.97, min(s * 0.22, 0.22))


def _theme_from_base(base_hex: str, mode: str, font: str, chart_palette: str) -> Theme:
    accent, gradient, bg = _derive_palette(base_hex, mode)
    return Theme(accent_color=accent, gradient_accent=gradient, background_color=bg,
                 font_heading=font, font_body=font, chart_palette=chart_palette)


# Each theme is generated from ONE base hue (+ light/dark mode) so its accent, impact
# gradient, and background are guaranteed coherent. Chart palettes are single-hue ramps
# matching each theme's family, never multi-hue (viridis/tab10/Set2) which broke coherence.
# Every theme is generated from ONE base hue on a LIGHT (near-white) background — no dark/black
# backgrounds, per user preference. The names ending in "_dark" are kept for backward
# compatibility with existing decks, but they now render on a light background like the rest.
THEME_CATALOG: dict = {
    "corporate_blue": _theme_from_base("#1E6FD0", "light", "Calibri", "Blues"),
    "warm_editorial": _theme_from_base("#E0612F", "light", "Calibri", "Oranges"),
    "minimal_mono":   _theme_from_base("#475569", "light", "Calibri", "Greys"),
    "vibrant_teal":   _theme_from_base("#0D9488", "light", "Calibri", "Greens"),
    "royal_purple":   _theme_from_base("#7C3AED", "light", "Calibri", "Purples"),
    "crimson_report": _theme_from_base("#C2333A", "light", "Calibri", "Reds"),
    "forest_green":   _theme_from_base("#15803D", "light", "Calibri", "Greens"),
    "amber_gold":     _theme_from_base("#B7791F", "light", "Calibri", "YlOrBr"),
    # Formerly dark-background themes — now light too (names kept for compatibility):
    "modern_dark":    _theme_from_base("#3B82F6", "light", "Calibri", "Blues"),
    "emerald_dark":   _theme_from_base("#10B981", "light", "Calibri", "Greens"),
    "graphite_dark":  _theme_from_base("#64748B", "light", "Calibri", "Greys"),
}
DEFAULT_THEME = "corporate_blue"


class StyleOverride(BaseModel):
    """Optional, explicit visual choice for a single slide, layered over the active theme.
    Absent in the common case; every field here takes precedence over the theme default
    for that one slide only."""
    background: Optional[BackgroundField] = None
    header_bar: Optional[bool] = None


class CoverSlide(BaseModel):
    type: Literal["cover"] = "cover"
    title: str
    subtitle: Optional[str] = ""
    date: Optional[str] = ""
    notes: Optional[str] = None


class ContentImageSlide(BaseModel):
    type: Literal["content_image"] = "content_image"
    title: str
    text: str
    image_id: str
    style_override: Optional[StyleOverride] = None
    notes: Optional[str] = None


class TwoColumnSide(BaseModel):
    title: str
    text: str


class TwoColumnSlide(BaseModel):
    type: Literal["two_column"] = "two_column"
    title: str
    left: TwoColumnSide
    right: TwoColumnSide
    style_override: Optional[StyleOverride] = None
    notes: Optional[str] = None


class SectionDividerSlide(BaseModel):
    type: Literal["section_divider"] = "section_divider"
    title: str
    subtitle: Optional[str] = ""
    notes: Optional[str] = None


class ContentTextSlide(BaseModel):
    type: Literal["content_text"] = "content_text"
    title: str
    text: str
    style_override: Optional[StyleOverride] = None
    notes: Optional[str] = None


class ChartData(BaseModel):
    """A chart expressed as one of four visualization intents instead of a low-level
    chart kind. The system maps each intent internally to the appropriate rendering
    in utils/charts.py and applies the active theme's chart palette automatically."""
    intent: Literal["comparison", "trend", "distribution", "part_of_whole"]
    title: Optional[str] = None
    categories: Optional[List[str]] = None
    values: Optional[List[float]] = None
    x: Optional[List[float]] = None
    y: Optional[List[float]] = None

    @model_validator(mode="after")
    def validate_chart(self):
        if self.intent in ("comparison", "part_of_whole"):
            if not self.categories or not self.values:
                raise ValueError(f"'{self.intent}' charts require 'categories' and 'values'.")
            if len(self.categories) != len(self.values):
                min_len = min(len(self.categories), len(self.values))
                self.categories = self.categories[:min_len]
                self.values = self.values[:min_len]
        elif self.intent == "trend":
            if not self.x or not self.y:
                raise ValueError("'trend' charts require 'x' and 'y'.")
            if len(self.x) != len(self.y):
                min_len = min(len(self.x), len(self.y))
                self.x = self.x[:min_len]
                self.y = self.y[:min_len]
        elif self.intent == "distribution":
            if not self.values:
                raise ValueError("'distribution' charts require 'values'.")
        return self


class TableData(BaseModel):
    headers: List[str]
    rows: List[List[str]]


class ContentMixedSlide(BaseModel):
    type: Literal["content_mixed"] = "content_mixed"
    title: str
    text: Optional[str] = None
    image_id: Optional[str] = None
    chart: Optional[ChartData] = None
    table: Optional[TableData] = None
    style_override: Optional[StyleOverride] = None
    notes: Optional[str] = None

    @model_validator(mode="after")
    def validate_contents(self):
        contents = [bool(self.image_id), bool(self.chart), bool(self.table)]
        if sum(contents) > 1:
            # AI assistants often set both image_id AND chart by mistake.
            # Keep only one, preferring: chart > table > image_id (richer wins).
            logger.warning(
                "content_mixed slide '%s' has %d visual element(s). Keeping only one.",
                self.title, sum(contents),
            )
            if self.chart and self.image_id:
                logger.warning("  -> Dropped image_id in favour of chart.")
                self.image_id = None
            elif self.chart and self.table:
                logger.warning("  -> Dropped table in favour of chart.")
                self.table = None
            elif self.table and self.image_id:
                logger.warning("  -> Dropped image_id in favour of table.")
                self.image_id = None
        elif sum(contents) == 0:
            # Instead of raising an error, log a warning and set a default
            logger.warning(
                "content_mixed slide '%s' has no visual element. Defaulting to text-only.",
                self.title,
            )
        return self


class TimelineEvent(BaseModel):
    fecha: str
    titulo: str
    emoji: Optional[str] = None


class TimelineSlide(BaseModel):
    type: Literal["timeline"] = "timeline"
    title: str
    items: List[TimelineEvent]
    active_index: Optional[int] = None
    # Two visual styles: "horizontal" (axis across the slide, best for 3-5 events) and
    # "vertical" (top-to-bottom rail, best for 4-6 longer events). In "vertical" style an
    # optional `image_id` or `text` renders on the LEFT and shifts the rail to the right.
    style: Literal["horizontal", "vertical"] = "horizontal"
    image_id: Optional[str] = None
    text: Optional[str] = None
    style_override: Optional[StyleOverride] = None
    notes: Optional[str] = None

    @model_validator(mode="after")
    def validate_items(self):
        if not self.items or len(self.items) < 2:
            raise ValueError("timeline requires at least two items.")
        if self.active_index is not None and not (0 <= self.active_index < len(self.items)):
            raise ValueError("timeline.active_index must refer to a valid item index.")
        return self


class ContentLatexSlide(BaseModel):
    type: Literal["content_latex"] = "content_latex"
    layout: Literal["full", "split"] = "split"
    title: str
    text: Optional[str] = None
    latex_lines: List[str] = Field(
        ...,
        description=(
            "Ordered list of mathtext strings to render as a stacked image. "
            "Wrap equations in $...$, e.g. '$E = mc^2$'. "
            "Use '$\\bullet\\;$' prefix for bullet-style items, "
            "e.g. '$\\bullet\\;$ Step 1: $F = ma$'."
        ),
    )
    image_id: Optional[str] = None
    style_override: Optional[StyleOverride] = None
    notes: Optional[str] = None


class StatHighlightSlide(BaseModel):
    """Spotlights a single key figure/metric with the same high-impact gradient/accent
    treatment as cover and section_divider — for the one number that should land hard."""
    type: Literal["stat_highlight"] = "stat_highlight"
    value: str
    label: str
    supporting_text: Optional[str] = None
    style_override: Optional[StyleOverride] = None
    notes: Optional[str] = None


Slide = Annotated[
    CoverSlide | ContentImageSlide | ContentMixedSlide | ContentLatexSlide | ContentTextSlide | TimelineSlide | TwoColumnSlide | SectionDividerSlide | StatHighlightSlide,
    Field(discriminator="type"),
]


class PresentationDefinition(BaseModel):
    theme: str = DEFAULT_THEME
    slides: List[Slide]

    @model_validator(mode="after")
    def validate_theme(self):
        if self.theme not in THEME_CATALOG:
            logger.warning(
                "Unrecognized theme '%s'; falling back to default theme '%s'.",
                self.theme, DEFAULT_THEME,
            )
            self.theme = DEFAULT_THEME
        return self


def hex_to_rgb(value: str) -> RGBColor:
    v = value.lstrip("#")
    return RGBColor(int(v[0:2], 16), int(v[2:4], 16), int(v[4:6], 16))


def _resolve_background(theme: "Theme", background: str):
    """Resolve background token or hex to (bg_rgb, txt_color).
    Text color is chosen automatically for contrast via relative luminance."""
    if background == "accent_color":
        bg_hex = theme.accent_color
    elif background == "background_color":
        bg_hex = theme.background_color
    else:
        bg_hex = background
    bg_rgb = hex_to_rgb(bg_hex)
    r, g, b = int(bg_rgb[0]), int(bg_rgb[1]), int(bg_rgb[2])
    luminance = (0.2126 * r + 0.7152 * g + 0.0722 * b) / 255.0
    txt_color = RGBColor(255, 255, 255) if luminance < 0.5 else RGBColor(30, 30, 30)
    return bg_rgb, txt_color


def _resolve_slide_background(theme: "Theme", override: Optional[StyleOverride]):
    """Resolve a slide's effective background, honoring an optional style_override."""
    token = override.background if (override and override.background) else "background_color"
    return _resolve_background(theme, token)


def _resolve_header_bar(override: Optional[StyleOverride], default: bool = True) -> bool:
    """Resolve a slide's effective header_bar flag, honoring an optional style_override."""
    if override is not None and override.header_bar is not None:
        return override.header_bar
    return default


import re as _re

# ─────────────────── YAML pre-processor for AI-generated YAML ───────────────────
# AI assistants often produce subtly invalid YAML. This function repairs the
# most common mistakes before the real parser sees the text, so that minor
# formatting errors don't cause total failure.


def _preprocess_yaml_text(raw: str) -> str:
    """Repair common AI-generated YAML mistakes in-place.

    Covers the most frequent errors seen across dozens of real AI outputs:
      - Duplicated keys: `title: "title": "value"` -> `title: "value"`
      - Orphan lines: `-The Math` (dash + word, no colon)
      - `text: "text": "value"` -> `text: "value"`
      - Duplicate keys inside inline dicts: `{x:[1,2], x:x}` -> `{x:[1,2]}`
      - Unquoted bare words in `y: [col1, col2]` lists
      - `file_name` / `OUTPUT` bare words at end of YAML
      - Missing space after dash in list items: `-type:` -> `- type:`
      - Completely empty slides: `- type: content_mixed` followed by `- type:`
    """
    text = raw
    fix_count = [0]  # mutable list so nested functions can increment

    # 1. Strip BOM and stray control chars (except \t \n \r)
    text = _re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f\x7f]', '', text)

    # 2. Remove orphan lines like "-The Math" (dash-start, no colon, not a proper list item)
    cleaned = []
    for line in text.splitlines():
        s = line.strip()
        if s.startswith('-') and ':' not in s and not s.startswith('- ') and not s.startswith('-{') and not s.startswith('-\n'):
            if _re.match(r'^-\w', s):
                fix_count[0] += 1
                continue
        cleaned.append(line)
    text = '\n'.join(cleaned)

    # 3. Fix duplicated keys: `title: "title": "value"` -> `title: "value"`
    text = _re.sub(r'(\w[\w_]*):\s*"\1"\s*:\s*', r'\1: ', text)
    # inline dict variant: {key: "key": value} -> {key: value}
    text = _re.sub(r'\{\s*(\w[\w_]*)\s*:\s*"\1"\s*:\s*', r'{\1: ', text)

    # 4. Fix `"text": "text": "..."` -> `"text": "..."`
    text = _re.sub(r'(?:["\x27])?text(?:["\x27])?\s*:\s*["\x27]text["\x27]\s*:\s*', 'text: ', text)

    # 5. Deduplicate last-wins keys inside inline data dicts: {x:[1,2], x:x} -> {x:x}
    def _dedup_inline_dict(m):
        block = m.group(0)
        fix_count[0] += 1
        brace_start = block.index('{')
        brace_end = block.rindex('}')
        inner = block[brace_start + 1:brace_end]
        pairs = _re.split(r',(?=(?:[^\[\]]*\[[^\[\]]*\])*[^\[\]]*$)', inner)
        seen = {}
        for pair in pairs:
            pair = pair.strip()
            if ':' not in pair:
                continue
            k, v = pair.split(':', 1)
            seen[k.strip()] = v.strip()
        deduped = ', '.join(f'{k}: {v}' for k, v in seen.items())
        return block[:brace_start + 1] + deduped + '}'

    # Match data: {...} where the dict may contain nested lists
    text = _re.sub(r'\bdata\s*:\s*\{[^}]+\}', _dedup_inline_dict, text)

    # 6. Fix unquoted bare words in y: [word1, word2] lists where they look like
    #    column name references (should be quoted strings for proper YAML parsing
    #    when they aren't valid YAML booleans/null/numbers)
    def _quote_bare_list_items(m):
        prefix = m.group(1)  # the key like "y" or "columns"
        inner = m.group(2)
        items = []
        for item in _re.split(r',\s*', inner):
            item = item.strip()
            # If it looks like a bare word (not number, not quoted, not bool/null)
            if (_re.match(r'^[a-zA-Z_]\w*$', item)
                    and item.lower() not in ('true', 'false', 'null', '~', 'yes', 'no', 'on', 'off')):
                items.append(f'"{item}"')
            else:
                items.append(item)
        return f'{prefix}: [{", ".join(items)}]'

    text = _re.sub(r'\b(y|columns)\s*:\s*\[([^\]]+)\]', _quote_bare_list_items, text)

    # 7. Remove trailing bare words that are not valid YAML (e.g. "file_name" or "OUTPUT" at end)
    text = _re.sub(r'\n\s*(file_name|OUTPUT|output)\s*(\n|$)', '\n', text, flags=_re.IGNORECASE)

    # 8. Fix missing space after dash in list items: `-type:` -> `- type:`
    text = _re.sub(r'^(\s*)-(\w+)', r'\1- \2', text, flags=_re.MULTILINE)

    # 9. Remove completely empty slides: `- type: content_mixed` followed by `- type: ...`
    #    (model sometimes emits a slide header with no data)
    text = _re.sub(
        r'(- type: content_mixed\s*\n)(?=- type:)',
        '',
        text,
    )

    if fix_count[0] > 0:
        logger.info("YAML pre-processor applied %d fix(es)", fix_count[0])

    return text


# Map common LaTeX commands to Unicode equivalents used in sanitize_slide_text
_LATEX_CMD_TO_UNICODE = {
    r'\alpha': 'α', r'\beta': 'β', r'\gamma': 'γ', r'\delta': 'δ',
    r'\epsilon': 'ε', r'\zeta': 'ζ', r'\eta': 'η', r'\theta': 'θ',
    r'\iota': 'ι', r'\kappa': 'κ', r'\lambda': 'λ', r'\mu': 'μ',
    r'\nu': 'ν', r'\xi': 'ξ', r'\pi': 'π', r'\rho': 'ρ',
    r'\sigma': 'σ', r'\tau': 'τ', r'\upsilon': 'υ', r'\phi': 'φ',
    r'\chi': 'χ', r'\psi': 'ψ', r'\omega': 'ω',
    r'\Gamma': 'Γ', r'\Delta': 'Δ', r'\Theta': 'Θ', r'\Lambda': 'Λ',
    r'\Xi': 'Ξ', r'\Pi': 'Π', r'\Sigma': 'Σ', r'\Upsilon': 'Υ',
    r'\Phi': 'Φ', r'\Psi': 'Ψ', r'\Omega': 'Ω',
    r'\nabla': '∇', r'\partial': '∂', r'\infty': '∞',
    r'\cdot': '·', r'\times': '×', r'\sqrt': '√',
    r'\leq': '≤', r'\geq': '≥', r'\neq': '≠', r'\approx': '≈',
    r'\in': '∈', r'\notin': '∉', r'\subset': '⊂',
    r'\sum': 'Σ', r'\prod': 'Π', r'\int': '∫',
    r'\rightarrow': '→', r'\leftarrow': '←',
    r'\Rightarrow': '⇒', r'\Leftarrow': '⇐',
    r'\leftrightarrow': '↔', r'\hat': '', r'\vec': '', r'\bar': '',
    r'\tilde': '', r'\frac': '',
}

_INLINE_BULLET_SEP = _re.compile(r'(?<=[.;:])\s+-\s+(?=\S)')


def _split_inline_dash_bullets(text: str) -> str:
    """Convert inline dash-separated clauses like 'A. - B. - C.' into separate
    bullet lines. AI-generated text sometimes uses ' - ' as a clause separator
    instead of real newlines, which otherwise renders as literal dash characters
    embedded in a paragraph rather than as a bulleted list."""
    lines = text.split('\n')
    out_lines = []
    for line in lines:
        stripped = line.strip()
        if not stripped or _re.match(r'^(\s*[-*•]|\s*\d+\.)\s+', stripped):
            out_lines.append(line)
            continue
        parts = _INLINE_BULLET_SEP.split(stripped)
        if len(parts) > 1:
            out_lines.append('- ' + parts[0])
            out_lines.extend('- ' + p for p in parts[1:])
        else:
            out_lines.append(line)
    return '\n'.join(out_lines)


def sanitize_slide_text(text: str, preserve_markdown: bool = False) -> str:
    """Strip markdown formatting and inline LaTeX from plain slide text.
    If preserve_markdown=True, keep bold/italic/code markers so the PowerPoint
    markdown renderer can apply formatting while still normalizing headings
    and math delimiters.
    """
    if not text:
        return text
    # Literal \n sequences — convert before $ processing so \nabla isn't split
    text = text.replace('\\n', '\n')
    # $$...$$ block equations — strip delimiters, keep inner content
    text = _re.sub(r'\$\$(.+?)\$\$', lambda m: m.group(1).strip(), text, flags=_re.DOTALL)
    # $...$ inline equations — strip delimiters, keep inner content
    text = _re.sub(r'\$(.+?)\$', lambda m: m.group(1).strip(), text)
    if not preserve_markdown:
        # Bold+italic ***...*** or ______
        text = _re.sub(r'[*_]{3}(.+?)[*_]{3}', r'\1', text)
        # Bold **...** or __...__
        text = _re.sub(r'[*_]{2}(.+?)[*_]{2}', r'\1', text)
        # Italic *...* or _..._
        text = _re.sub(r'[*_](.+?)[*_]', r'\1', text)
        # Inline code `...`
        text = _re.sub(r'`(.+?)`', r'\1', text)
    # ATX headings: leading #+
    text = _re.sub(r'^#{1,6}\s+', '', text, flags=_re.MULTILINE)
    # Replace known LaTeX commands with Unicode equivalents
    for cmd, uni in _LATEX_CMD_TO_UNICODE.items():
        text = text.replace(cmd, uni)
    # \cmd{arg} → arg  (e.g. \hat{y} → y)
    text = _re.sub(r'\\[a-zA-Z]+\{([^}]*)\}', r'\1', text)
    # Remaining bare \commands → strip
    text = _re.sub(r'\\[a-zA-Z]+', '', text)
    text = _split_inline_dash_bullets(text)
    return text


_UNSUPPORTED_MATHTEXT = (
    r'\bigg', r'\Bigg', r'\big', r'\Big',
    r'\left', r'\right',
    r'\mkern', r'\mspace', r'\hspace', r'\vspace',
    r'\mathrm', r'\mathit', r'\mathbf',
)

def _sanitize_latex_line(line: str) -> str:
    r"""Normalize a mathtext line for matplotlib rendering.

    - Converts $$...$$ delimiters to $...$ (matplotlib uses single-dollar mathtext)
    - Collapses over-escaped backslashes (\\\\cmd -> \\cmd -> \cmd in mathtext)
    - Removes sizing/spacing commands not supported by matplotlib mathtext
    - Ensures the line has $...$ delimiters if it looks like math
    """
    if not line:
        return line
    # Convert $$...$$ to $...$ (model often outputs display-math delimiters)
    line = line.replace('$$', '$')
    # Translate markdown-style bold/italic into mathtext commands
    line = _re.sub(r'\*\*(.+?)\*\*', r'\\mathbf{\1}', line)
    line = _re.sub(r'\*(.+?)\*', r'\\mathit{\1}', line)
    # Collapse double-backslash to single (over-escaped by model)
    while '\\\\' in line:
        line = line.replace('\\\\', '\\')
    # Remove unsupported commands, keeping just their argument if parenthesised
    for cmd in _UNSUPPORTED_MATHTEXT:
        line = line.replace(cmd, '')
    return line


def set_slide_background(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _apply_gradient_background(slide, color1: RGBColor, color2: RGBColor, angle: float = 45.0):
    """Two-stop diagonal gradient background, used only on "impact" slides
    (cover, section_divider, stat_highlight) — never on ordinary content slides."""
    fill = slide.background.fill
    fill.gradient()
    stops = fill.gradient_stops
    stops[0].color.rgb = color1
    stops[0].position = 0.0
    stops[1].color.rgb = color2
    stops[1].position = 1.0
    fill.gradient_angle = angle


def _blend_color(c1: RGBColor, c2: RGBColor, t: float) -> RGBColor:
    """Linear-interpolate between two colors. t=0 -> c1, t=1 -> c2."""
    return RGBColor(
        int(c1[0] + (c2[0] - c1[0]) * t),
        int(c1[1] + (c2[1] - c1[1]) * t),
        int(c1[2] + (c2[2] - c1[2]) * t),
    )


def _scale_lightness(rgb: RGBColor, factor: float) -> RGBColor:
    """Scale a color's lightness by `factor`, keeping its hue and saturation. This varies
    the DEPTH of a color (lighter/deeper) without changing its family — a blue stays blue."""
    import colorsys
    h, l, s = colorsys.rgb_to_hls(rgb[0] / 255, rgb[1] / 255, rgb[2] / 255)
    l = max(0.16, min(0.90, l * factor))
    r, g, b = colorsys.hls_to_rgb(h, l, s)
    return RGBColor(int(r * 255), int(g * 255), int(b * 255))


def _impact_gradient(theme: "Theme", variant: int = 0):
    """Return one (color1, color2, angle) gradient that is IDENTICAL on every impact slide
    (cover, section_divider, stat_highlight) so the deck's "moments" share exactly the same
    look. Earlier we varied lightness per slide, but that read as inconsistent ("not the
    same blue"); uniformity is what looks harmonious. `variant` is accepted for call-site
    compatibility but intentionally ignored.

    Dark themes use a deeper gradient so impact slides stay rich (cohesive with the charcoal
    content slides) rather than turning into a glaringly bright panel.

    Returns (color1, color2, angle, text_color). The text color (white or near-black) is the
    one that contrasts better, and the gradient is nudged deeper/lighter until BOTH stops clear
    a WCAG ratio of 4.5 with it — so impact-slide titles are always legible, for any theme."""
    accent = hex_to_rgb(theme.accent_color)
    grad = hex_to_rgb(theme.gradient_accent)
    bg = hex_to_rgb(theme.background_color)
    white, black = RGBColor(255, 255, 255), RGBColor(20, 20, 20)
    dark_theme = _relative_luminance(bg) < 0.2

    def min_contrast(t):
        return min(_wcag_contrast(t, c1), _wcag_contrast(t, c2))

    if dark_theme:
        # Dark themes: a deep, rich gradient with WHITE text, cohesive with the charcoal
        # content slides (never a glaringly bright panel that clashes with the rest).
        c1, c2 = _scale_lightness(accent, 0.72), _scale_lightness(grad, 0.62)
        use_white = True
    else:
        # Light themes: keep the vivid accent gradient; pick whichever text contrasts better
        # (e.g. white on deep blue, dark on bright gold).
        c1, c2 = accent, grad
        use_white = min_contrast(white) >= min_contrast(black)

    text = white if use_white else black
    factor = 0.9 if use_white else 1.12  # darken for white text, lighten for dark text
    for _ in range(14):
        if min_contrast(text) >= 4.5:
            break
        c1, c2 = _scale_lightness(c1, factor), _scale_lightness(c2, factor)
    return c1, c2, 55.0, text


def _relative_luminance(c: RGBColor) -> float:
    """WCAG relative luminance (gamma-corrected)."""
    def _lin(v: float) -> float:
        v = int(v) / 255.0
        return v / 12.92 if v <= 0.03928 else ((v + 0.055) / 1.055) ** 2.4
    return 0.2126 * _lin(c[0]) + 0.7152 * _lin(c[1]) + 0.0722 * _lin(c[2])


def _wcag_contrast(c1: RGBColor, c2: RGBColor) -> float:
    """WCAG contrast ratio between two colors (1.0 = none, 21.0 = black/white)."""
    l1, l2 = _relative_luminance(c1), _relative_luminance(c2)
    hi, lo = max(l1, l2), min(l1, l2)
    return (hi + 0.05) / (lo + 0.05)


def _header_fill(accent_rgb: RGBColor) -> RGBColor:
    """Fill color for header bars / column-title chips. White title text sits on top, so the
    fill is deepened until it clears a comfortable WCAG contrast (>=4.5) against white — for
    ANY hue. Greens/teals look bright for their lightness, so a fixed factor isn't enough;
    iterating on measured contrast guarantees legibility across every theme."""
    fill = accent_rgb
    white = RGBColor(255, 255, 255)
    for _ in range(10):
        if _wcag_contrast(fill, white) >= 4.5:
            break
        fill = _scale_lightness(fill, 0.82)
    return fill


def _contrast_text_color(*colors: RGBColor) -> RGBColor:
    """Contrasting text color (white or near-black) for one or more background colors,
    based on their average relative luminance."""
    def _luminance(c: RGBColor) -> float:
        r, g, b = int(c[0]), int(c[1]), int(c[2])
        return (0.2126 * r + 0.7152 * g + 0.0722 * b) / 255.0
    avg = sum(_luminance(c) for c in colors) / len(colors)
    return RGBColor(255, 255, 255) if avg < 0.55 else RGBColor(30, 30, 30)


def _parse_inline_markdown(text: str):
    """Parse inline markdown and return list of (segment_text, bold, italic) tuples.
    Supports: **bold**, *italic*, ***bold+italic***, `code` (treated as bold mono).
    """
    import re
    pattern = re.compile(
        r'(\*\*\*(?P<bolditalic>.+?)\*\*\*)'
        r'|(\*\*(?P<bold>.+?)\*\*)'
        r'|(\*(?P<italic>.+?)\*)'
        r'|(__(?P<bold2>.+?)__)'
        r'|(_(?P<italic2>.+?)_)'
        r'|(`(?P<code>.+?)`)',
        re.DOTALL,
    )
    segments = []
    last = 0
    for m in pattern.finditer(text):
        if m.start() > last:
            segments.append((text[last:m.start()], False, False))
        if m.group('bolditalic'):
            segments.append((m.group('bolditalic'), True, True))
        elif m.group('bold') or m.group('bold2'):
            t = m.group('bold') or m.group('bold2')
            segments.append((t, True, False))
        elif m.group('italic') or m.group('italic2'):
            t = m.group('italic') or m.group('italic2')
            segments.append((t, False, True))
        elif m.group('code'):
            segments.append((m.group('code'), True, False))
        last = m.end()
    if last < len(text):
        segments.append((text[last:], False, False))
    return segments if segments else [(text, False, False)]


def _add_run_to_paragraph(p, text, font_name, font_size, bold, italic, color):
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return run


def _parse_md_table(table_lines: list) -> Optional['TableData']:
    """Parse a list of markdown table lines (header, separator, rows) into TableData."""
    if len(table_lines) < 2:
        return None

    def split_row(line):
        return [c.strip() for c in line.strip().strip('|').split('|')]

    headers = split_row(table_lines[0])
    # table_lines[1] is the separator — skip it
    rows = [split_row(line) for line in table_lines[2:] if line.strip() and '|' in line]
    n = len(headers)
    normalized = []
    for row in rows:
        if len(row) < n:
            row = row + [''] * (n - len(row))
        elif len(row) > n:
            row = row[:n]
        normalized.append(row)
    if not headers:
        return None
    return TableData(headers=headers, rows=normalized)


def _split_text_and_tables(text: str):
    """Split markdown text into segments: ('text', str) or ('table', TableData).
    Detects standard pipe-table syntax: header row | separator row | data rows.
    """
    lines = text.splitlines()
    segments = []
    current_lines: list = []
    i = 0
    while i < len(lines):
        line = lines[i]
        # Pipe table: current line has | and next line is all dashes/colons/pipes
        if ('|' in line and line.strip().startswith('|')
                and i + 1 < len(lines)
                and _re.match(r'^\s*\|[\s\-:|]+\|\s*$', lines[i + 1])):
            if current_lines:
                combined = '\n'.join(current_lines).strip()
                if combined:
                    segments.append(('text', combined))
                current_lines = []
            table_lines = [line]
            i += 1
            while i < len(lines) and '|' in lines[i]:
                table_lines.append(lines[i])
                i += 1
            td = _parse_md_table(table_lines)
            if td:
                segments.append(('table', td))
        else:
            current_lines.append(line)
            i += 1
    if current_lines:
        combined = '\n'.join(current_lines).strip()
        if combined:
            segments.append(('text', combined))
    return segments if segments else [('text', text)]


def _make_paragraph_bullet_pointed(p):
    pPr = p._p.get_or_add_pPr()
    pPr.set('marL', '342900')
    pPr.set('indent', '-171450')
    buFont = OxmlElement('a:buFont')
    buFont.set('typeface', 'Arial')
    buFont.set('panose', '020B0604020202020204')
    buFont.set('pitchFamily', '34')
    buFont.set('charset', '0')
    pPr.append(buFont)
    buChar = OxmlElement('a:buChar')
    buChar.set('char', '•')
    pPr.append(buChar)


def _fill_textbox_markdown(tf, text, font_name, font_size, bold, italic, color, align):
    """Fill an existing text frame with markdown-formatted lines."""
    lines = text.replace('\\n', '\n').splitlines()
    first_paragraph = True
    for line in lines:
        stripped = line.rstrip()
        is_bullet = bool(_re.match(r'^(\s*[-*•]\s+|\s*\d+\.\s+)', stripped))
        if is_bullet:
            m = _re.match(r'^(\s*)([-*•]|\d+\.)\s+', stripped)
            if m:
                stripped = stripped[m.end():]
            if first_paragraph:
                p = tf.paragraphs[0]
                first_paragraph = False
            else:
                p = tf.add_paragraph()
            _make_paragraph_bullet_pointed(p)
            p.alignment = PP_ALIGN.LEFT
            p.level = 0
        else:
            if first_paragraph:
                p = tf.paragraphs[0]
                first_paragraph = False
            else:
                p = tf.add_paragraph()
            p.alignment = align
        # Breathing room so body text doesn't read as a cramped block hugging the top.
        p.line_spacing = 1.25
        p.space_after = Pt(10)
        for seg_text, seg_bold, seg_italic in _parse_inline_markdown(stripped):
            _add_run_to_paragraph(p, seg_text, font_name, font_size,
                                  bold or seg_bold, italic or seg_italic, color)


def _force_bullet_text(text: str) -> str:
    lines = text.replace('\\n', '\n').splitlines()
    normalized = []
    for line in lines:
        stripped = line.strip()
        if not stripped:
            normalized.append('')
            continue
        if _re.match(r'^(\s*[-*•]|\s*\d+\.)\s+', stripped):
            normalized.append(stripped)
        else:
            normalized.append(f'- {stripped}')
    return '\n'.join(normalized)


def _autofit_font_size(text: str, base_size: int, width_emu: int, height_emu: int) -> int:
    """Estimate the largest font size (down to a floor) at which `text` fits in a box of the
    given size, so dense slides don't spill past the slide edge. Coarse on purpose — it errs
    toward shrinking; PowerPoint's normAutofit then refines the exact scale on open."""
    import math
    width_in = max(width_emu / 914400.0, 0.5)
    height_in = max(height_emu / 914400.0, 0.5)
    lines = (text or "").replace("\\n", "\n").splitlines()
    n_paras = max(1, len([l for l in lines if l.strip()]))
    size = base_size
    while size > 12:
        chars_per_line = max(1, int(width_in * 142 / size))  # ~avg Calibri glyph width
        wrapped = 0.0
        for l in lines:
            s = l.strip()
            wrapped += math.ceil(len(s) / chars_per_line) if s else 0.5
        total_in = wrapped * (size * 1.25) / 72.0 + n_paras * (10 / 72.0)
        if total_in <= height_in:
            break
        size -= 1
    return size


def add_text_box(
    slide,
    left, top, width, height,
    text,
    font_name, font_size,
    bold=False, italic=False,
    color=RGBColor(0, 0, 0),
    align=PP_ALIGN.CENTER,
    word_wrap=True,
    markdown=False,
    text_style: Literal["prose", "bullets"] = "prose",
    accent_rgb: RGBColor = None,
    bg_rgb: RGBColor = None,
    vertical_center: bool = False,
    autofit: bool = False,
):
    """Add a text box. markdown=True parses bold/italic/bullets.
    text_style='bullets' forces every non-empty line into a bullet item.
    Embedded markdown pipe tables are automatically detected and rendered as real PPTX tables.
    vertical_center=True anchors short content to the middle of the box instead of the top,
    so sparse text doesn't leave a large empty area below it.
    autofit=True shrinks the font to fit the box ("Shrink text on overflow") so dense slides
    never spill text past the slide edge; combined here with a content-length pre-scale.
    """
    if text and text_style == "bullets":
        text = _force_bullet_text(text)

    # Pre-shrink long body text so it fits even in renderers that don't apply PowerPoint's
    # autofit lazily. Heuristic on raw text length; PowerPoint's normAutofit refines further.
    if autofit and text:
        font_size = _autofit_font_size(text, font_size, width, height)

    def _apply_autofit(tf):
        if autofit:
            tf.word_wrap = True
            tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

    if not markdown or not text:
        shape = slide.shapes.add_textbox(left, top, width, height)
        shape.text_frame.word_wrap = word_wrap
        if vertical_center:
            shape.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        shape.text_frame.clear()
        p = shape.text_frame.paragraphs[0]
        p.alignment = align
        _add_run_to_paragraph(p, text or '', font_name, font_size, bold, italic, color)
        _apply_autofit(shape.text_frame)
        return shape

    # Detect embedded markdown tables
    segments = _split_text_and_tables(text.replace('\\n', '\n'))
    has_table = any(seg_type == 'table' for seg_type, _ in segments)

    if not has_table:
        shape = slide.shapes.add_textbox(left, top, width, height)
        shape.text_frame.word_wrap = word_wrap
        if vertical_center:
            shape.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        shape.text_frame.clear()
        _fill_textbox_markdown(shape.text_frame, text, font_name, font_size,
                               bold, italic, color, align)
        _apply_autofit(shape.text_frame)
        return shape

    # Mixed text + table: stack segments vertically within the given bounds
    gap = Inches(0.12)
    n_text = sum(1 for t, _ in segments if t == 'text')
    tbl_heights = []
    for seg_type, seg_content in segments:
        if seg_type == 'table':
            ideal = (len(seg_content.rows) + 1) * Inches(0.40)
            tbl_heights.append(min(ideal, Inches(2.5)))
    total_tbl_h = sum(tbl_heights)
    total_gaps = (len(segments) - 1) * gap
    remaining_text_h = max(height - total_tbl_h - total_gaps, Inches(0.4) * max(n_text, 1))
    text_h_each = int(remaining_text_h // max(n_text, 1))

    y = top
    tbl_idx = 0
    last_shape = None
    for seg_type, seg_content in segments:
        remaining = top + height - y
        if remaining < Inches(0.2):
            break
        if seg_type == 'text':
            seg_h = min(text_h_each, remaining)
            shape = slide.shapes.add_textbox(left, y, width, seg_h)
            shape.text_frame.word_wrap = word_wrap
            shape.text_frame.clear()
            _fill_textbox_markdown(shape.text_frame, seg_content, font_name, font_size,
                                   bold, italic, color, align)
            last_shape = shape
            y += seg_h + gap
        else:
            seg_h = min(tbl_heights[tbl_idx], remaining)
            tbl_idx += 1
            add_table(slide, seg_content, left, y, width, seg_h, font_name,
                      accent_rgb=accent_rgb, bg_rgb=bg_rgb, txt_color=color)
            y += seg_h + gap
    return last_shape


def add_header_bar(slide, prs, title, font_name, accent_rgb, bar_height=Inches(0.8)):
    fill = _header_fill(accent_rgb)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, bar_height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = fill
    bar.line.fill.background()

    text_box = slide.shapes.add_textbox(0, 0, prs.slide_width, bar_height)
    text_box.fill.background()
    text_box.line.fill.background()
    tf = text_box.text_frame
    tf.margin_top = Inches(0.15)
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.word_wrap = False
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = title
    run.font.name = font_name
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.color.rgb = _contrast_text_color(fill)
    return bar_height


def place_image_centered(slide, img_source, left, top, max_width, max_height,
                          valign: str = "center", halign: str = "center"):
    if hasattr(img_source, "seek"):
        img_source.seek(0)
    with PILImage.open(img_source) as im:
        img_w, img_h = im.size

    ratio = img_w / img_h
    if (max_width / max_height) >= ratio:
        final_h = max_height
        final_w = int(final_h * ratio)
    else:
        final_w = max_width
        final_h = int(final_w / ratio)

    if halign == "left":
        offset_x = 0
    elif halign == "right":
        offset_x = int(max_width - final_w)
    else:
        offset_x = int((max_width - final_w) / 2)

    if valign == "top":
        offset_y = 0
    elif valign == "bottom":
        offset_y = int(max_height - final_h)
    else:
        offset_y = int((max_height - final_h) / 2)

    if hasattr(img_source, "seek"):
        img_source.seek(0)
    picture = slide.shapes.add_picture(
        img_source,
        left + offset_x,
        top + offset_y,
        width=final_w,
        height=final_h,
    )
    return picture


def add_rect(slide, x, y, w, h, fill_color, line_color=None, line_w_pt=0):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        int(round(x)), int(round(y)), int(round(w)), int(round(h))
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_w_pt)
    else:
        shape.line.fill.background()
    return shape


def add_oval_shape(slide, x, y, w, h, fill_color, line_color=None, line_w_pt=1.5):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        int(round(x)), int(round(y)), int(round(w)), int(round(h))
    )
    if fill_color is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_w_pt)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_line_shape(slide, x1, y1, x2, y2, color, width_pt=1.5, dash=False):
    connector = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        int(round(x1)), int(round(y1)), int(round(x2)), int(round(y2))
    )
    connector.line.color.rgb = color
    connector.line.width = Pt(width_pt)
    if dash:
        connector.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    return connector


_INTENT_TO_CHART_KIND = {
    "comparison": "bar",
    "trend": "line",
    "distribution": "hist",
    "part_of_whole": "pie",
}


def _render_chart_image(chart_def: ChartData, palette: str, width, height):
    """Render one of the four visualization intents using the existing seaborn/matplotlib
    backend in utils/charts.py. The theme's chart_palette is always applied; the LLM
    never selects a low-level chart kind or palette directly."""
    kind = _INTENT_TO_CHART_KIND[chart_def.intent]

    # Use the slide's chart title if given, else NO title (a single space suppresses the
    # backend's ugly auto-generated "Barras: value por category" default). Axis labels are
    # blanked too — the category names already appear on the axis, so "category"/"value"
    # are noise that looks unprofessional.
    kwargs = {
        "save_path": None,
        "palette": palette,
        "title": chart_def.title if chart_def.title else " ",
        "xlabel": " ",
        "ylabel": " ",
    }

    if chart_def.intent in ("comparison", "part_of_whole"):
        df = pd.DataFrame({"category": chart_def.categories, "value": chart_def.values})
        kwargs["x"] = "category"
        kwargs["y"] = "value"
    elif chart_def.intent == "trend":
        df = pd.DataFrame({"x": chart_def.x, "y": chart_def.y})
        kwargs["x"] = "x"
        kwargs["y"] = "y"
    else:  # distribution
        df = pd.DataFrame({"value": chart_def.values})
        kwargs["x"] = "value"

    from matplotlib import pyplot as plt

    buf = BytesIO()
    kwargs["save_path"] = buf
    result = chart(kind, df, **kwargs)

    if hasattr(result, "figure"):
        plt.close(result.figure)
    elif isinstance(result, tuple):
        first = result[0]
        if hasattr(first, "figure"):
            plt.close(first.figure)
        elif isinstance(first, plt.Figure):
            plt.close(first)
        else:
            plt.close('all')
    elif isinstance(result, plt.Figure):
        plt.close(result)
    else:
        plt.close('all')

    buf.seek(0)
    return buf


def add_chart(slide, chart_def: ChartData, palette: str, left, top, width, height,
              accent_rgb: RGBColor = None, bg_rgb: RGBColor = None,
              txt_color: RGBColor = None):
    """Render charts using the matplotlib backend and insert the image into PowerPoint.
    The chart is placed on a white rounded card with a thin themed border so it reads as
    an intentionally placed element rather than a stray white rectangle, especially on
    dark or strongly colored themes (the matplotlib backend always renders on white)."""
    try:
        chart_image = _render_chart_image(chart_def, palette, width, height)
        pad = Inches(0.12)
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            int(left - pad), int(top - pad), int(width + 2 * pad), int(height + 2 * pad),
        )
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(255, 255, 255)
        card.line.color.rgb = accent_rgb or RGBColor(200, 200, 200)
        card.line.width = Pt(0.75)
        card.shadow.inherit = False
        try:
            card.adjustments[0] = 0.04
        except (IndexError, ValueError):
            pass
        place_image_centered(slide, chart_image, left, top, width, height)
    except Exception as e:
        logger.warning("Chart rendering failed: %s. Adding placeholder text instead.", e)
        add_text_box(slide, left, top, width, height,
                     f"[Chart could not be rendered: {e}]",
                     "Calibri", 12, color=txt_color or RGBColor(180, 180, 180),
                     align=PP_ALIGN.LEFT)
    return

def add_table(slide, table_def: TableData, left, top, width, height, font_name,
              font_size: int = 11, accent_rgb: RGBColor = None,
              bg_rgb: RGBColor = None, txt_color: RGBColor = None,
              vertical_center: bool = False):
    """Add a styled table: accent-colored header row, alternating row backgrounds.
    vertical_center=True centers the table within the allotted `height` so it lines up with
    vertically-centered body text beside it (instead of hugging the top of its column)."""
    _accent = accent_rgb or RGBColor(0x0D, 0x94, 0x88)
    _bg = bg_rgb or RGBColor(0xF0, 0xFD, 0xFA)
    _txt = txt_color or RGBColor(30, 30, 30)
    # Slightly tinted alternate row color
    alt_rgb = RGBColor(
        min(255, int(_bg[0] * 0.93 + _accent[0] * 0.07)),
        min(255, int(_bg[1] * 0.93 + _accent[1] * 0.07)),
        min(255, int(_bg[2] * 0.93 + _accent[2] * 0.07)),
    )
    rows = len(table_def.rows) + 1
    cols = len(table_def.headers)
    row_height = Inches(0.40)
    actual_height = min(rows * row_height, height)
    if vertical_center and actual_height < height:
        top = int(top + (height - actual_height) / 2)
    tbl_shape = slide.shapes.add_table(rows, cols, left, top, width, actual_height)
    tbl = tbl_shape.table
    # Header row
    for ci, header in enumerate(table_def.headers):
        cell = tbl.cell(0, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = _accent
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.runs[0] if p.runs else p.add_run()
        run.text = str(header)
        run.font.name = font_name
        run.font.size = Pt(font_size)
        run.font.bold = True
        run.font.color.rgb = RGBColor(255, 255, 255)
    # Data rows
    for ri, row_data in enumerate(table_def.rows, start=1):
        row_bg = alt_rgb if ri % 2 == 0 else _bg
        for ci in range(cols):
            cell = tbl.cell(ri, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = row_bg
            cell_text = str(row_data[ci]) if ci < len(row_data) else ''
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.runs[0] if p.runs else p.add_run()
            run.text = cell_text
            run.font.name = font_name
            run.font.size = Pt(font_size)
            run.font.color.rgb = _txt
    return tbl_shape


def add_vertical_separator(slide, x, top, bottom, accent_rgb):
    sep = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x, top, x, bottom)
    sep.line.width = Pt(1.5)
    sep.line.color.rgb = accent_rgb


MARGIN = Inches(0.6)
BAR_H = Inches(0.8)


def _add_notes(slide, text: Optional[str]) -> None:
    """Write presenter notes to the slide's notes pane."""
    if not text:
        return
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = text


def _render_latex_to_image(
    latex_lines: list,
    bg_rgb: RGBColor,
    txt_rgb: RGBColor,
    dpi: int = 300,
) -> BytesIO:
    """Render a list of mathtext strings to a transparent PNG image in memory."""
    import matplotlib
    matplotlib.use('Agg')
    import matplotlib.pyplot as plt

    fg = (txt_rgb[0] / 255, txt_rgb[1] / 255, txt_rgb[2] / 255)
    sanitized_lines = [_sanitize_latex_line(raw_line).strip() for raw_line in latex_lines] or ['']
    # Lines with tall multi-part notation (summations, fractions, integrals, roots)
    # need more vertical room than a single-height row, or adjacent lines visually
    # touch each other once rendered.
    weights = [
        1.5 if any(marker in line for marker in (r'\sum', r'\prod', r'\int', r'\frac', r'\sqrt'))
        else 1.0
        for line in sanitized_lines
    ]
    total_weight = sum(weights)
    fig_h = max(1.5, total_weight * 0.95)
    fig, ax = plt.subplots(figsize=(10, fig_h), dpi=dpi)
    fig.patch.set_alpha(0)
    ax.patch.set_alpha(0)
    ax.set_xlim(0, 1)
    ax.set_ylim(0, 1)
    ax.axis('off')
    # matplotlib parses mathtext LAZILY at draw/savefig time, not when ax.text() is
    # called — so a single invalid line (e.g. \textbf{...}, an unbalanced brace, or an
    # unsupported command) would raise inside savefig and blank the ENTIRE equation image.
    # To stop one bad line from wiping the whole slide, validate each line with the
    # mathtext parser up front and render any line that fails as plain (non-math) text.
    from matplotlib import mathtext as _mathtext
    _math_parser = _mathtext.MathTextParser('path')

    def _to_plain(s: str) -> str:
        s = _re.sub(r'\$', '', s)
        for _cmd, _uni in _LATEX_CMD_TO_UNICODE.items():
            s = s.replace(_cmd, _uni)
        s = _re.sub(r'\\[a-zA-Z]+\{([^}]*)\}', r'\1', s)   # \text{x} -> x
        s = _re.sub(r'\\[a-zA-Z]+', '', s)                  # bare \cmd -> (drop)
        return s.replace('{', '').replace('}', '').strip()

    def _math_ok(s: str) -> bool:
        try:
            _math_parser.parse(s)
            return True
        except Exception:
            return False

    def _draw(plain_only: bool):
        cumulative = 0.0
        for line, weight in zip(sanitized_lines, weights):
            candidate = line
            if candidate and not (candidate.startswith('$') and candidate.endswith('$')):
                candidate = f'${candidate}$'
            y = 1.0 - (cumulative + weight / 2) / total_weight
            cumulative += weight
            if not plain_only and candidate and _math_ok(candidate):
                ax.text(0.02, y, candidate, transform=ax.transAxes, fontsize=18,
                        color=fg, verticalalignment='center', horizontalalignment='left')
            else:
                ax.text(0.02, y, _to_plain(line), transform=ax.transAxes, fontsize=18,
                        color=fg, verticalalignment='center', horizontalalignment='left',
                        parse_math=False)

    _draw(plain_only=False)
    buf = BytesIO()
    try:
        fig.savefig(buf, format='png', dpi=dpi, bbox_inches='tight',
                    transparent=True, edgecolor='none')
    except Exception:
        # Last-resort safety net: re-render every line as plain text so the slide is
        # never blank, even if a line slipped past validation.
        ax.clear()
        ax.set_xlim(0, 1)
        ax.set_ylim(0, 1)
        ax.axis('off')
        _draw(plain_only=True)
        buf = BytesIO()
        fig.savefig(buf, format='png', dpi=dpi, bbox_inches='tight',
                    transparent=True, edgecolor='none')
    plt.close(fig)
    buf.seek(0)
    return buf


def _resolve_image(image_registry: dict, image_id: str):
    image = image_registry.get(image_id)
    if image is None:
        raise FileNotFoundError(f"image_id '{image_id}' not found in registry.")
    if hasattr(image, "seek"):
        image.seek(0)
    return image


def build_cover(prs, theme: Theme, data: CoverSlide, image_registry: dict, variant: int = 0):
    """High-impact cover: a clean two-tone gradient whose hue varies per impact slide,
    with content vertically and horizontally centered. No decorative shapes."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    c1, c2, angle, txt = _impact_gradient(theme, variant)
    _apply_gradient_background(slide, c1, c2, angle)
    W, H = prs.slide_width, prs.slide_height
    secondary = _blend_color(txt, RGBColor(128, 128, 128), 0.25)
    inner_x = MARGIN + Inches(0.4)
    inner_w = W - (MARGIN + Inches(0.4)) * 2

    # Compute total block height for vertical centering
    TITLE_H = Inches(1.6)
    SUB_H = Inches(0.55) if data.subtitle else Inches(0)
    DATE_H = Inches(0.35) if data.date else Inches(0)
    GAP_TS = Inches(0.28) if data.subtitle else Inches(0)
    GAP_SD = Inches(0.18) if data.date else Inches(0)
    block_h = TITLE_H + GAP_TS + SUB_H + GAP_SD + DATE_H
    block_top = int((H - block_h) / 2)

    add_text_box(slide, inner_x, block_top, inner_w, TITLE_H,
                 data.title, theme.font_heading, 46,
                 bold=True, color=txt, align=PP_ALIGN.CENTER)
    y = block_top + TITLE_H + GAP_TS
    if data.subtitle:
        add_text_box(slide, inner_x, y, inner_w, SUB_H,
                     data.subtitle, theme.font_body, 22,
                     italic=True, color=secondary, align=PP_ALIGN.CENTER)
        y += SUB_H + GAP_SD
    if data.date:
        add_text_box(slide, inner_x, y, inner_w, DATE_H,
                     data.date, theme.font_body, 14,
                     color=secondary, align=PP_ALIGN.CENTER)
    _add_notes(slide, data.notes)


def _add_accent_title_bar(slide, prs, title, font_name, accent_rgb, txt_color, header_bar, bg_rgb):
    """Unified title rendering. Returns content_top (Inches position after title area).
    - header_bar=True: full-width colored bar with white text
    - header_bar=False: plain title + thin accent underline rule
    """
    W = prs.slide_width
    if header_bar:
        return add_header_bar(slide, prs, title, font_name, accent_rgb)
    else:
        title_top = Inches(0.35)
        title_h = Inches(0.65)
        add_text_box(slide, MARGIN, title_top, W - MARGIN * 2, title_h,
                     title, font_name, 26, bold=True, color=txt_color, align=PP_ALIGN.LEFT)
        # thin accent underline
        rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN, title_top + title_h, W - MARGIN * 2, Inches(0.04))
        rule.fill.solid()
        rule.fill.fore_color.rgb = accent_rgb
        rule.line.fill.background()
        return title_top + title_h + Inches(0.1)


def build_content_image(prs, theme: Theme, data: ContentImageSlide, image_registry: dict, variant: int = 0):
    """Left-text / Right-image layout."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg_rgb, txt_color = _resolve_slide_background(theme, data.style_override)
    accent = hex_to_rgb(theme.accent_color)
    set_slide_background(slide, bg_rgb)
    W, H = prs.slide_width, prs.slide_height
    header_bar = _resolve_header_bar(data.style_override, default=True)
    content_top = _add_accent_title_bar(slide, prs, data.title, theme.font_heading, accent, txt_color, header_bar, bg_rgb)
    content_top += Inches(0.3)
    content_h = H - content_top - MARGIN
    gutter = Inches(0.4)
    text_col_w = int((W - MARGIN * 2 - gutter) * 0.42)
    img_col_w = int((W - MARGIN * 2 - gutter) * 0.58)
    add_text_box(slide, MARGIN, content_top, text_col_w, content_h,
                 sanitize_slide_text(data.text, preserve_markdown=True), theme.font_body, 17,
                 color=txt_color, align=PP_ALIGN.LEFT, word_wrap=True, markdown=True,
                 vertical_center=True, autofit=True)
    img_x = MARGIN + text_col_w + gutter
    img_source = _resolve_image(image_registry, data.image_id)
    place_image_centered(slide, img_source, img_x, content_top + Inches(0.1), img_col_w, content_h - Inches(0.2))
    _add_notes(slide, data.notes)


def build_two_column(prs, theme: Theme, data: TwoColumnSlide, image_registry: dict, variant: int = 0):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg_rgb, txt_color = _resolve_slide_background(theme, data.style_override)
    accent = hex_to_rgb(theme.accent_color)
    set_slide_background(slide, bg_rgb)
    W, H = prs.slide_width, prs.slide_height
    header_bar = _resolve_header_bar(data.style_override, default=False)
    content_top = _add_accent_title_bar(slide, prs, data.title, theme.font_heading, accent, txt_color, header_bar, bg_rgb)
    content_top += Inches(0.3)  # breathing room below main title
    content_h = H - content_top - MARGIN
    gutter = Inches(0.5)  # wider gap provides visual separation without a line
    col_w = int((W - MARGIN * 2 - gutter) / 2)
    _build_column(slide, MARGIN, content_top, col_w, content_h,
                  sanitize_slide_text(data.left.title), data.left.text,
                  theme.font_heading, theme.font_body, txt_color, accent)
    _build_column(slide, MARGIN + col_w + gutter, content_top, col_w, content_h,
                  sanitize_slide_text(data.right.title), data.right.text,
                  theme.font_heading, theme.font_body, txt_color, accent)
    _add_notes(slide, data.notes)


def _build_column(slide, left, top, width, height,
                  title, text, font_heading, font_body, txt_color, accent):
    """Render a single column: accent-colored header rectangle + body text."""
    title_h = Inches(0.52)
    # Accent-colored header rectangle defining the column title zone. Deepen a bright accent
    # so the white title text stays legible (matches the slide header bar treatment).
    header_fill = _header_fill(accent)
    header_rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, title_h)
    header_rect.fill.solid()
    header_rect.fill.fore_color.rgb = header_fill
    header_rect.line.fill.background()
    # Title text over the colored header (contrast-aware against the chip fill)
    add_text_box(slide, left + Inches(0.1), top + Inches(0.05),
                 width - Inches(0.12), title_h - Inches(0.08),
                 title, font_heading, 15,
                 bold=True, color=_contrast_text_color(header_fill), align=PP_ALIGN.LEFT)
    # Body text below. Top-aligned (NOT vertically centered) so the two columns of a
    # comparison line up row-for-row even when they hold different amounts of text.
    body_top = top + title_h + Inches(0.22)
    body_h = height - title_h - Inches(0.22)
    add_text_box(slide, left, body_top, width, body_h,
                 sanitize_slide_text(text, preserve_markdown=True), font_body, 15,
                 color=txt_color, align=PP_ALIGN.LEFT, word_wrap=True, markdown=True,
                 vertical_center=False, autofit=True)


def build_content_mixed(prs, theme: Theme, data: ContentMixedSlide, image_registry: dict, variant: int = 0):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg_rgb, txt_color = _resolve_slide_background(theme, data.style_override)
    accent = hex_to_rgb(theme.accent_color)
    set_slide_background(slide, bg_rgb)
    W, H = prs.slide_width, prs.slide_height
    header_bar = _resolve_header_bar(data.style_override, default=True)
    content_top = _add_accent_title_bar(slide, prs, data.title, theme.font_heading, accent, txt_color, header_bar, bg_rgb)
    content_top += Inches(0.3)
    gutter = Inches(0.4)
    total_inner_w = W - MARGIN * 2 - gutter
    text_col_w = int(total_inner_w * 0.38)
    right_col_w = total_inner_w - text_col_w
    add_text_box(slide, MARGIN, content_top, text_col_w, H - content_top - MARGIN,
                 sanitize_slide_text(data.text or '', preserve_markdown=True), theme.font_body, 17,
                 color=txt_color, align=PP_ALIGN.LEFT, word_wrap=True, markdown=True,
                 accent_rgb=accent, bg_rgb=bg_rgb, vertical_center=True, autofit=True)
    right_x = MARGIN + text_col_w + gutter
    if data.image_id:
        img_source = _resolve_image(image_registry, data.image_id)
        img_top = content_top + Inches(0.08)
        img_height = H - content_top - MARGIN - Inches(0.16)
        place_image_centered(slide, img_source, right_x, img_top, right_col_w, img_height)
    elif data.chart:
        add_chart(slide, data.chart, theme.chart_palette, right_x, content_top, right_col_w, H - content_top - MARGIN,
                  accent_rgb=accent, bg_rgb=bg_rgb, txt_color=txt_color)
    elif data.table:
        add_table(slide, data.table, right_x, content_top + Inches(0.1), right_col_w,
                  H - content_top - MARGIN - Inches(0.1), theme.font_body,
                  accent_rgb=accent, bg_rgb=bg_rgb, txt_color=txt_color, vertical_center=True)
    _add_notes(slide, data.notes)


def build_timeline(prs, theme: Theme, data: TimelineSlide, image_registry: dict, variant: int = 0):
    if data.style == "vertical":
        return _build_timeline_vertical(prs, theme, data, image_registry)
    return _build_timeline_horizontal(prs, theme, data)


def _timeline_line_color(txt_color) -> RGBColor:
    """Subtle connector color that stays visible on both light and dark backgrounds."""
    is_dark_bg = int(txt_color[0]) > 128  # white text => dark slide background
    return RGBColor(0x3C, 0x4A, 0x60) if is_dark_bg else RGBColor(0xCB, 0xD5, 0xE1)


def _draw_timeline_node(slide, cx, cy, accent, bg_rgb, active: bool):
    """Draw a clean, professional marker: a small solid accent dot, with a thin hollow
    accent ring around the active one. No oversized emoji bubbles (those read as childish)."""
    dot_r = Inches(0.13) if active else Inches(0.085)
    if active:
        ring_r = Inches(0.23)
        # Hollow ring: fill matches the slide background so only the outline shows.
        add_oval_shape(slide, int(cx - ring_r), int(cy - ring_r), ring_r * 2, ring_r * 2,
                       fill_color=bg_rgb, line_color=accent, line_w_pt=1.75)
    add_oval_shape(slide, int(cx - dot_r), int(cy - dot_r), dot_r * 2, dot_r * 2,
                   fill_color=accent, line_color=None)


def _build_timeline_horizontal(prs, theme: Theme, data: TimelineSlide):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg_rgb, txt_color = _resolve_slide_background(theme, data.style_override)
    accent = hex_to_rgb(theme.accent_color)
    line_rgb = _timeline_line_color(txt_color)
    set_slide_background(slide, bg_rgb)
    W, H = prs.slide_width, prs.slide_height
    header_bar = _resolve_header_bar(data.style_override, default=True)
    content_top = _add_accent_title_bar(slide, prs, data.title, theme.font_heading, accent, txt_color, header_bar, bg_rgb)
    content_top += Inches(0.3)

    items = data.items
    n_items = len(items)
    active_idx = data.active_index if data.active_index is not None else -1
    # Inset the axis so first/last labels have room and nodes never touch the slide edge.
    inset = Inches(1.1)
    axis_left = MARGIN + inset
    axis_right = W - MARGIN - inset
    step = (axis_right - axis_left) / max(n_items - 1, 1)

    line_y = int(round(H * 0.50))
    # Thin refined axis.
    add_rect(slide, axis_left, line_y - int(Inches(0.0125)), int(axis_right - axis_left), Inches(0.025),
             fill_color=line_rgb)

    text_box_width = Inches(2.3)
    for i, item in enumerate(items):
        x = int(round(axis_left + i * step))
        active = (i == active_idx)
        _draw_timeline_node(slide, x, line_y, accent, bg_rgb, active)

        is_top = (i % 2 == 0)
        text_left = int(round(x - text_box_width / 2))
        text_left = max(int(MARGIN), min(int(W - MARGIN - text_box_width), text_left))
        gap = Inches(0.40)
        if is_top:
            date_top = int(round(line_y - gap - Inches(0.78)))
        else:
            date_top = int(round(line_y + gap))

        add_text_box(slide, text_left, date_top, text_box_width, Inches(0.26),
                     item.fecha, theme.font_body, 11, color=accent, bold=True, align=PP_ALIGN.CENTER)
        add_text_box(slide, text_left, date_top + Inches(0.28), text_box_width, Inches(0.52),
                     item.titulo, theme.font_body, 14, color=txt_color, bold=True,
                     align=PP_ALIGN.CENTER, word_wrap=True)

    _add_notes(slide, data.notes)


def _build_timeline_vertical(prs, theme: Theme, data: TimelineSlide, image_registry: dict):
    """Vertical timeline: a top-to-bottom rail with stacked nodes, each event's date + title
    to its right. With an optional `image_id`/`text`, that content fills the LEFT side and the
    rail shifts to the RIGHT half (a clean two-pane layout). Suits 4-6 events / longer titles."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg_rgb, txt_color = _resolve_slide_background(theme, data.style_override)
    accent = hex_to_rgb(theme.accent_color)
    line_rgb = _timeline_line_color(txt_color)
    set_slide_background(slide, bg_rgb)
    W, H = prs.slide_width, prs.slide_height
    header_bar = _resolve_header_bar(data.style_override, default=True)
    content_top = _add_accent_title_bar(slide, prs, data.title, theme.font_heading, accent, txt_color, header_bar, bg_rgb)
    content_top += Inches(0.25)
    content_h = H - content_top - MARGIN

    has_side = bool(data.image_id or data.text)
    if has_side:
        # Left pane: image or text. Right pane: the timeline rail + events.
        gutter = Inches(0.5)
        left_w = int((W - MARGIN * 2 - gutter) * 0.46)
        rail_zone_left = MARGIN + left_w + gutter
        if data.image_id:
            img_source = _resolve_image(image_registry, data.image_id)
            place_image_centered(slide, img_source, MARGIN, content_top, left_w, content_h)
        elif data.text:
            add_text_box(slide, MARGIN, content_top, left_w, content_h,
                         sanitize_slide_text(data.text, preserve_markdown=True), theme.font_body, 16,
                         color=txt_color, align=PP_ALIGN.LEFT, word_wrap=True, markdown=True,
                         accent_rgb=accent, bg_rgb=bg_rgb, vertical_center=True)
    else:
        rail_zone_left = MARGIN

    items = data.items
    n_items = len(items)
    active_idx = data.active_index if data.active_index is not None else -1
    row_h = content_h / n_items
    rail_x = rail_zone_left + Inches(0.30)
    centers = [int(round(content_top + row_h * i + row_h / 2)) for i in range(n_items)]

    # Continuous rail behind the nodes, from first to last center.
    rail_w = Inches(0.03)
    add_rect(slide, int(rail_x - rail_w / 2), centers[0], rail_w, centers[-1] - centers[0],
             fill_color=line_rgb)

    text_x = int(rail_x + Inches(0.45))
    text_w = W - text_x - MARGIN

    for i, item in enumerate(items):
        cy = centers[i]
        active = (i == active_idx)
        _draw_timeline_node(slide, rail_x, cy, accent, bg_rgb, active)

        block_h = Inches(0.78)
        block_top = int(cy - block_h / 2)
        add_text_box(slide, text_x, block_top, text_w, Inches(0.26),
                     item.fecha, theme.font_body, 11, color=accent, bold=True, align=PP_ALIGN.LEFT)
        add_text_box(slide, text_x, block_top + Inches(0.28), text_w, Inches(0.50),
                     item.titulo, theme.font_body, 15, color=txt_color, bold=True,
                     align=PP_ALIGN.LEFT, word_wrap=True)

    _add_notes(slide, data.notes)


def build_section_divider(prs, theme: Theme, data: SectionDividerSlide, image_registry: dict, variant: int = 0):
    """Section divider: a clean two-tone gradient whose hue varies per impact slide,
    centered title and subtitle. No decorative shapes."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    c1, c2, angle, txt = _impact_gradient(theme, variant)
    _apply_gradient_background(slide, c1, c2, angle)
    W, H = prs.slide_width, prs.slide_height
    inner_w = W - MARGIN * 2
    secondary = _blend_color(txt, RGBColor(128, 128, 128), 0.25)
    TITLE_H = Inches(1.1)
    SUB_H = Inches(0.65) if data.subtitle else Inches(0)
    GAP = Inches(0.25) if data.subtitle else Inches(0)
    block_h = TITLE_H + GAP + SUB_H
    block_top = int((H - block_h) / 2)
    add_text_box(slide, MARGIN, block_top, inner_w, TITLE_H,
                 data.title, theme.font_heading, 40,
                 bold=True, color=txt, align=PP_ALIGN.CENTER)
    if data.subtitle:
        add_text_box(slide, MARGIN, block_top + TITLE_H + GAP, inner_w, SUB_H,
                     data.subtitle, theme.font_body, 22,
                     italic=True, color=secondary, align=PP_ALIGN.CENTER)
    _add_notes(slide, data.notes)


def build_stat_highlight(prs, theme: Theme, data: StatHighlightSlide, image_registry: dict, variant: int = 0):
    """Spotlight a single key figure WITHOUT looking like a chapter opener. Unlike
    cover/section_divider (full-bleed gradient), this sits on the normal light content
    background and puts the figure inside a colored KPI "card". That still lands hard
    (a big number in a bold accent panel) but reads clearly as a highlight WITHIN the
    current section, not the start of a new one."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    W, H = prs.slide_width, prs.slide_height
    bg_rgb, txt_color = _resolve_slide_background(theme, data.style_override)
    accent = hex_to_rgb(theme.accent_color)
    set_slide_background(slide, bg_rgb)

    card_fill = _header_fill(accent)              # deep enough for white text
    card_txt = _contrast_text_color(card_fill)
    secondary = _blend_color(txt_color, RGBColor(128, 128, 128), 0.35)

    CARD_W = int((W - MARGIN * 2) * 0.62)
    CARD_H = Inches(2.5)
    LABEL_H = Inches(0.6)
    SUPPORT_H = Inches(0.5) if data.supporting_text else Inches(0)
    GAP_CARD_LABEL = Inches(0.35)
    GAP_LABEL_SUPPORT = Inches(0.15)
    block_h = CARD_H + GAP_CARD_LABEL + LABEL_H + (GAP_LABEL_SUPPORT + SUPPORT_H if data.supporting_text else Inches(0))
    block_top = int((H - block_h) / 2)
    card_left = int((W - CARD_W) / 2)

    # The KPI card: a rounded accent panel that frames the figure.
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, card_left, block_top, CARD_W, CARD_H)
    card.fill.solid()
    card.fill.fore_color.rgb = card_fill
    card.line.fill.background()
    card.shadow.inherit = False
    try:
        card.adjustments[0] = 0.06
    except (IndexError, ValueError):
        pass

    # Scale the figure to "land hard" but fit the card; autofit shrinks long values further.
    value_size = 170 if len(data.value) <= 4 else (130 if len(data.value) <= 8 else 96)
    add_text_box(slide, card_left + Inches(0.2), block_top, CARD_W - Inches(0.4), CARD_H,
                 data.value, theme.font_heading, value_size,
                 bold=True, color=card_txt, align=PP_ALIGN.CENTER, word_wrap=True,
                 vertical_center=True, autofit=True)

    y = block_top + CARD_H + GAP_CARD_LABEL
    add_text_box(slide, MARGIN, y, W - MARGIN * 2, LABEL_H,
                 data.label, theme.font_body, 24,
                 bold=True, color=txt_color, align=PP_ALIGN.CENTER, word_wrap=True)
    if data.supporting_text:
        y += LABEL_H + GAP_LABEL_SUPPORT
        add_text_box(slide, MARGIN, y, W - MARGIN * 2, SUPPORT_H,
                     sanitize_slide_text(data.supporting_text), theme.font_body, 14,
                     color=secondary, align=PP_ALIGN.CENTER, word_wrap=True)
    _add_notes(slide, data.notes)


def build_content_text(prs, theme: Theme, data: ContentTextSlide, image_registry: dict, variant: int = 0):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg_rgb, txt_color = _resolve_slide_background(theme, data.style_override)
    accent = hex_to_rgb(theme.accent_color)
    set_slide_background(slide, bg_rgb)
    W, H = prs.slide_width, prs.slide_height
    header_bar = _resolve_header_bar(data.style_override, default=True)
    content_top = _add_accent_title_bar(slide, prs, data.title, theme.font_heading, accent, txt_color, header_bar, bg_rgb)
    content_top += Inches(0.3)  # breathing room below title
    content_h = H - content_top - MARGIN
    add_text_box(slide, MARGIN, content_top, W - MARGIN * 2, content_h,
                 sanitize_slide_text(data.text, preserve_markdown=True), theme.font_body, 19,
                 color=txt_color, align=PP_ALIGN.LEFT, word_wrap=True, markdown=True,
                 accent_rgb=accent, bg_rgb=bg_rgb, vertical_center=True, autofit=True)
    _add_notes(slide, data.notes)


def build_content_latex(prs, theme: Theme, data: ContentLatexSlide, image_registry: dict, variant: int = 0):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg_rgb, txt_color = _resolve_slide_background(theme, data.style_override)
    accent = hex_to_rgb(theme.accent_color)
    set_slide_background(slide, bg_rgb)
    W, H = prs.slide_width, prs.slide_height
    header_bar = _resolve_header_bar(data.style_override, default=True)
    content_top = _add_accent_title_bar(slide, prs, data.title, theme.font_heading, accent, txt_color, header_bar, bg_rgb)
    content_top += Inches(0.3)
    content_h = H - content_top - MARGIN
    if data.layout == 'split' and data.image_id:
        # Split layout: equations image on the left, a normal image on the right
        gutter = Inches(0.4)
        col_w = int((W - MARGIN * 2 - gutter) / 2)
        latex_img = _render_latex_to_image(data.latex_lines, bg_rgb, txt_color)
        place_image_centered(slide, latex_img, MARGIN, content_top, col_w, content_h, valign="center")
        img_source = _resolve_image(image_registry, data.image_id)
        place_image_centered(slide, img_source, MARGIN + col_w + gutter, content_top, col_w, content_h, valign="center")
    elif data.layout == 'split' and data.text:
        # Split layout: descriptive text on the left, equations image on the right
        gutter = Inches(0.4)
        col_w = int((W - MARGIN * 2 - gutter) / 2)
        add_text_box(slide, MARGIN, content_top, col_w, content_h,
                     sanitize_slide_text(data.text, preserve_markdown=True), theme.font_body, 16,
                     color=txt_color, align=PP_ALIGN.LEFT, word_wrap=True, markdown=True,
                     vertical_center=True, autofit=True)
        latex_img = _render_latex_to_image(data.latex_lines, bg_rgb, txt_color)
        place_image_centered(slide, latex_img, MARGIN + col_w + gutter, content_top, col_w, content_h, valign="center")
    else:
        # Full layout: optional intro text above, then full-width equations image
        if data.text:
            text_h = Inches(0.65)
            add_text_box(slide, MARGIN, content_top, W - MARGIN * 2, text_h,
                         sanitize_slide_text(data.text, preserve_markdown=True), theme.font_body, 14,
                         color=txt_color, align=PP_ALIGN.LEFT, word_wrap=True, markdown=True)
            content_top += text_h + Inches(0.15)
            content_h = H - content_top - MARGIN
        latex_img = _render_latex_to_image(data.latex_lines, bg_rgb, txt_color)
        place_image_centered(slide, latex_img, MARGIN, content_top, W - MARGIN * 2, content_h, valign="center")
    _add_notes(slide, data.notes)


BUILDERS = {
    "cover":            build_cover,
    "content_image":    build_content_image,
    "content_mixed":    build_content_mixed,
    "content_latex":    build_content_latex,
    "content_text":     build_content_text,
    "timeline":         build_timeline,
    "two_column":       build_two_column,
    "section_divider":  build_section_divider,
    "stat_highlight":   build_stat_highlight,
}


def _strip_control_chars(obj):
    """Recursively remove control characters (except \\t \\n \\r) from all strings in a
    parsed YAML structure, so YAML escape sequences like \\b/\\f can't leak garbage
    control bytes into slide text or titles."""
    if isinstance(obj, str):
        return _re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f\x7f]', '', obj)
    if isinstance(obj, dict):
        return {k: _strip_control_chars(v) for k, v in obj.items()}
    if isinstance(obj, list):
        return [_strip_control_chars(v) for v in obj]
    return obj


def create_presentation_from_yaml(
    yaml_text: str,
    output_buffer,
    image_registry: dict = None,
) -> BytesIO:
    image_registry = image_registry or {}

    # Strip YAML-illegal control characters (\x00-\x08, \x0b-\x1f except \t\n\r)
    yaml_text = _re.sub(r'[\x00-\x08\x0b-\x0c\x0e-\x1f\x7f]', '', yaml_text)

    # Run the AI pre-processor to fix common AI generation mistakes
    yaml_text = _preprocess_yaml_text(yaml_text)

    try:
        raw = yaml.safe_load(yaml_text)
    except yaml.YAMLError as exc:
        raise ValueError(f"YAML parse error: {exc}") from exc

    if not isinstance(raw, dict):
        raise ValueError("YAML must be a mapping at the top level.")

    # Strip control characters introduced by YAML escape sequences during parsing
    # (e.g. a double-quoted "\beta" becomes backspace+"eta"). Without this, control
    # chars survive into titles — which are NOT run through sanitize_slide_text — and
    # serialize as garbage like "_x0008_eta_1" in the PPTX.
    raw = _strip_control_chars(raw)

    # The simplified schema has no 'global' block; a legacy/retired 'global' key is
    # silently dropped (logged) rather than auto-migrated, per the full schema replacement.
    if 'global' in raw:
        logger.warning(
            "Ignoring legacy 'global' block - this tool now uses a single 'theme' field "
            "instead of global.accent_color/background_color/font_heading/font_body."
        )
        raw.pop('global', None)

    # Normalize: common alternate names for 'slides'
    if 'slides' not in raw:
        for alt in ('slide_list', 'slide', 'pages', 'content'):
            if alt in raw:
                raw['slides'] = raw.pop(alt)
                break

    presentation_def = PresentationDefinition.model_validate(raw)
    theme = THEME_CATALOG[presentation_def.theme]
    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)

    _IMPACT_TYPES = {"cover", "section_divider", "stat_highlight"}
    impact_variant = 0
    for i, slide_data in enumerate(presentation_def.slides):
        builder = BUILDERS.get(slide_data.type)
        if not builder:
            logger.warning("Slide %d: type '%s' not implemented, skipping.", i, slide_data.type)
            continue
        variant = 0
        if slide_data.type in _IMPACT_TYPES:
            variant = impact_variant
            impact_variant += 1
        try:
            builder(prs, theme, slide_data, image_registry, variant)
        except Exception as slide_e:
            slide_title = getattr(slide_data, 'title', f'slide {i}')
            logger.warning(
                "Slide %d ('%s') failed: %s. Skipping to keep the rest of the presentation intact.",
                i, slide_title, slide_e,
            )
            continue

    prs.save(output_buffer)
    output_buffer.seek(0)
    return output_buffer


def _extract_image_ids_from_yaml(parsed_yaml) -> list[str]:
    image_ids: list[str] = []

    if isinstance(parsed_yaml, dict):
        for key, value in parsed_yaml.items():
            if key == "image_id" and isinstance(value, str):
                image_ids.append(value)
            else:
                image_ids.extend(_extract_image_ids_from_yaml(value))
    elif isinstance(parsed_yaml, list):
        for item in parsed_yaml:
            image_ids.extend(_extract_image_ids_from_yaml(item))

    return image_ids


def generate_powerpoint_structured_yaml(
    document_yaml,
    file_name,
    request,
    URL,
    ENABLE_CREATE_KNOWLEDGE,
    knowledge_name,
):
    try:
        try:
            parsed_yaml = yaml.safe_load(document_yaml)
        except yaml.YAMLError as exc:
            return {"error": {"message": f"YAML parse error: {exc}"}}

        image_ids = list(dict.fromkeys(_extract_image_ids_from_yaml(parsed_yaml)))
        images = []

        if image_ids:
            logger.info(f"Received {len(image_ids)} image references for PPTX YAML generation.")

        for image_id in image_ids:
            image_file = download_file(URL, _get_bearer_token(request), image_id)
            if isinstance(image_file, dict) and "error" in image_file:
                return {"error": {"message": f"Error downloading image with ID {image_id}: {image_file['error']['message']}"}}
            image_file.seek(0)
            images.append(image_file)

        image_registry = {image_id: image_file for image_id, image_file in zip(image_ids, images)}
        buffer = BytesIO()
        buffer.name = f"{file_name}.pptx"

        try:
            create_presentation_from_yaml(document_yaml, buffer, image_registry)
        except Exception as exec_e:
            error_msg = str(exec_e)
            # Add helpful hints based on the error message
            hints = []
            if "YAML parse error" in error_msg:
                hints.append(
                    "Suggestion: Check for duplicated keys like `title: \"title\": \"value\"` "
                    "or unquoted bare words in lists like `y: [col1, col2]`. "
                    "Ensure every list item starts with `- ` (dash + space)."
                )
            elif "content_mixed" in error_msg and "image_id" in error_msg:
                hints.append(
                    "Suggestion: A content_mixed slide has both image_id AND chart/table. "
                    "Use ONLY ONE visual element per slide."
                )
            elif "charts require" in error_msg:
                hints.append(
                    "Suggestion: 'comparison'/'part_of_whole' charts need 'categories' + 'values'; "
                    "'trend' charts need 'x' + 'y'; 'distribution' charts need 'values'."
                )
            if hints:
                error_msg = error_msg + "\n\n" + "\n".join(hints)
            return {"error": {"message": f"Error generating PowerPoint from YAML: {error_msg}"}}

        buffer.seek(0)

        try:
            bearer_token = _get_bearer_token(request)
            if bearer_token:
                logger.info("=> Received authorization header!")
            else:
                logger.warning("=> Authorization header is missing.")
        except Exception:
            logger.exception("=> Error retrieving authorization header")
            bearer_token = None

        user_id = None
        if bearer_token:
            user_id = get_user_id(URL, bearer_token)
            if not user_id:
                logger.warning("=> Could not obtain user id from token; continuing without knowledge creation.")

        response, request_data = upload_file(
            url=URL,
            token=bearer_token,
            file_data=buffer,
            filename=file_name,
            file_type="pptx"
        )

        if "file_path_download" in response and ENABLE_CREATE_KNOWLEDGE and user_id:
            create_knowledge_status = create_knowledge(
                url=URL,
                token=bearer_token,
                file_id=request_data['id'],
                user_id=user_id,
                knowledge_name=knowledge_name
            )
            if create_knowledge_status:
                logger.info("=> Knowledge base updated successfully.")
            else:
                logger.error("=> Error creating or updating knowledge base")
        elif "file_path_download" in response and ENABLE_CREATE_KNOWLEDGE and not user_id:
            logger.warning("=> Skipping knowledge creation because user_id is unavailable.")
        elif "error" in response:
            logger.error("=> Error uploading the file.")
        else:
            logger.info("=> Skipping knowledge creation because ENABLE_CREATE_KNOWLEDGE is false")

        return response
    except Exception as e:
        logger.error("=> An unknown error occurred during .PPTX document generation from YAML.")
        return dumps(
            {
                "error": {
                    "message": str(e)
                }
            },
            indent=4,
            ensure_ascii=False
        )